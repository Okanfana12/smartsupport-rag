"""
src/ingestion/loader.py
═══════════════════════
Loader universel — SmartSupport RAG

Rôle :
    Détecter le format d'un fichier et charger son contenu
    en utilisant le loader approprié.
    Retourne toujours List[Document] — même interface pour tous les formats.

Formats supportés :
    PDF, DOCX, CSV, JSON, HTML, TXT, MD, PPTX, EML

Choix techniques :
    - JSON  → json natif Python (pas de jq requis)
    - EML   → email natif Python (pas de unstructured requis)
    - PPTX  → python-pptx natif (pas de unstructured requis)
    → Zéro dépendance lourde externe

Auteur : Oumou Kanfana
"""

import os
import json
import email as email_lib
from pathlib import Path
from typing import List

from langchain_core.documents import Document
from langchain_community.document_loaders import (
    PyPDFLoader,      # PDF  — extrait texte page par page
    Docx2txtLoader,   # DOCX — extrait texte des paragraphes Word
    CSVLoader,        # CSV  — une ligne = un Document
    BSHTMLLoader,     # HTML — retire les balises, garde le texte
    TextLoader,       # TXT  — lit le texte brut ligne par ligne
)


# =============================================================
# Mapping extension → loader LangChain
# JSON, EML, PPTX sont gérés séparément via des loaders natifs
# =============================================================
LOADER_MAP = {
    ".pdf":  PyPDFLoader,
    ".docx": Docx2txtLoader,
    ".csv":  CSVLoader,
    ".html": BSHTMLLoader,
    ".htm":  BSHTMLLoader,
    ".txt":  TextLoader,
    ".md":   TextLoader,
}

# Extensions supportées — inclut les formats natifs Python
SUPPORTED_EXTENSIONS = set(LOADER_MAP.keys()) | {".json", ".eml", ".pptx"}


# =============================================================
# Loaders natifs Python — sans dépendances lourdes
# =============================================================

def load_json(file_path: str) -> List[Document]:
    """
    Charge un fichier JSON sans dépendance externe.

    Stratégie :
        - JSON liste → 1 Document par élément
        - JSON dict  → 1 Document par clé de premier niveau
        - Autre      → 1 seul Document

    Args:
        file_path : chemin vers le fichier JSON

    Returns:
        List[Document] avec le contenu JSON converti en texte
    """
    with open(file_path, 'r', encoding='utf-8') as f:
        data = json.load(f)

    documents = []

    if isinstance(data, list):
        # JSON tableau → 1 Document par élément
        for i, item in enumerate(data):
            content = json.dumps(item, ensure_ascii=False, indent=2)
            documents.append(Document(
                page_content=content,
                metadata={"source": file_path, "index": i}
            ))

    elif isinstance(data, dict):
        # JSON objet → 1 Document par clé de premier niveau
        for key, value in data.items():
            content = f"{key}:\n{json.dumps(value, ensure_ascii=False, indent=2)}"
            documents.append(Document(
                page_content=content,
                metadata={"source": file_path, "key": key}
            ))

    else:
        documents.append(Document(
            page_content=str(data),
            metadata={"source": file_path}
        ))

    return documents


def load_eml(file_path: str) -> List[Document]:
    """
    Charge un fichier EML avec le module email natif Python.
    Pas de dépendance externe — pas de unstructured requis.

    Extrait : sujet, expéditeur, date, corps du message.

    Args:
        file_path : chemin vers le fichier EML

    Returns:
        List[Document] avec le contenu de l'email
    """
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
        msg = email_lib.message_from_file(f)

    # Extraction des en-têtes
    subject = msg.get('Subject', '')
    sender  = msg.get('From', '')
    date    = msg.get('Date', '')

    # Extraction du corps du message
    body = ""
    if msg.is_multipart():
        # Email multipart → chercher la partie texte
        for part in msg.walk():
            if part.get_content_type() == 'text/plain':
                payload = part.get_payload(decode=True)
                if payload:
                    body += payload.decode('utf-8', errors='ignore')
    else:
        payload = msg.get_payload(decode=True)
        if payload:
            body = payload.decode('utf-8', errors='ignore')
        else:
            # Payload non encodé — lire directement
            body = str(msg.get_payload())

    content = f"Sujet: {subject}\nDe: {sender}\nDate: {date}\n\n{body}"

    return [Document(
        page_content=content,
        metadata={"source": file_path}
    )]


def load_pptx(file_path: str) -> List[Document]:
    """
    Charge un fichier PPTX avec python-pptx natif.
    Pas de dépendance externe — pas de unstructured requis.

    Chaque slide devient un Document séparé.

    Args:
        file_path : chemin vers le fichier PPTX

    Returns:
        List[Document] — 1 Document par slide
    """
    from pptx import Presentation

    prs = Presentation(file_path)
    documents = []

    for i, slide in enumerate(prs.slides):
        # Extraction du texte de chaque slide
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())

        content = "\n".join(texts)

        if content:
            documents.append(Document(
                page_content=content,
                metadata={"source": file_path, "slide": i + 1}
            ))

    return documents


# =============================================================
# Fonctions principales
# =============================================================

def detect_extension(file_path: str) -> str:
    """
    Détecte l'extension d'un fichier en minuscules.

    Args:
        file_path : chemin vers le fichier

    Returns:
        Extension en minuscules — ex: ".pdf", ".csv"

    Raises:
        ValueError : si le fichier n'a pas d'extension
    """
    ext = Path(file_path).suffix.lower()

    if not ext:
        raise ValueError(
            f"Impossible de détecter le format : {file_path}\n"
            f"Assurez-vous que le fichier a une extension."
        )

    return ext


def get_loader(file_path: str, ext: str):
    """
    Retourne le loader LangChain approprié selon l'extension.

    JSON, EML, PPTX retournent None — traités par leurs
    fonctions dédiées dans load_file().

    Args:
        file_path : chemin vers le fichier
        ext       : extension du fichier

    Returns:
        Instance du loader LangChain ou None

    Raises:
        ValueError : si le format n'est pas supporté
    """
    # Formats natifs Python — traités séparément
    if ext in {".json", ".eml", ".pptx"}:
        return None

    loader_class = LOADER_MAP.get(ext)

    if loader_class is None:
        supported = ", ".join(sorted(SUPPORTED_EXTENSIONS))
        raise ValueError(
            f"Format non supporté : {ext}\n"
            f"Formats acceptés : {supported}"
        )

    return loader_class(file_path)


def add_metadata(documents: List[Document], file_path: str) -> List[Document]:
    """
    Enrichit chaque Document avec des métadonnées sur le fichier source.

    Args:
        documents : liste de Documents chargés
        file_path : chemin vers le fichier source

    Returns:
        Documents enrichis avec filename, file_path, file_type, doc_index
    """
    path = Path(file_path)

    for i, doc in enumerate(documents):
        doc.metadata["filename"]  = path.name
        doc.metadata["file_path"] = str(path)
        doc.metadata["file_type"] = path.suffix.lower()
        doc.metadata["doc_index"] = i

    return documents


def load_file(file_path: str) -> List[Document]:
    """
    Charge un fichier et retourne une liste de Documents LangChain.

    Fonction principale du module — un seul appel pour tous les formats.

    Args:
        file_path : chemin vers le fichier à charger

    Returns:
        List[Document] — même format quelle que soit la source

    Raises:
        FileNotFoundError : si le fichier n'existe pas
        ValueError        : si le format n'est pas supporté

    Exemple:
        >>> docs = load_file("data/pdf_files/FAQ_Support.pdf")
        >>> print(len(docs))
        >>> print(docs[0].page_content[:200])
        >>> print(docs[0].metadata)
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(
            f"Fichier introuvable : {file_path}"
        )

    ext = detect_extension(file_path)

    # ── JSON — loader natif Python ────────────────────────────
    if ext == ".json":
        documents = load_json(file_path)
        documents = add_metadata(documents, file_path)
        print(f"[Loader] ✅ {len(documents)} document(s) chargé(s) depuis '{Path(file_path).name}'")
        return documents

    # ── EML — loader natif Python ─────────────────────────────
    if ext == ".eml":
        documents = load_eml(file_path)
        documents = add_metadata(documents, file_path)
        print(f"[Loader] ✅ {len(documents)} document(s) chargé(s) depuis '{Path(file_path).name}'")
        return documents

    # ── PPTX — loader natif python-pptx ──────────────────────
    if ext == ".pptx":
        documents = load_pptx(file_path)
        documents = add_metadata(documents, file_path)
        print(f"[Loader] ✅ {len(documents)} document(s) chargé(s) depuis '{Path(file_path).name}'")
        return documents

    # ── Autres formats — loader LangChain ────────────────────
    loader = get_loader(file_path, ext)
    documents = loader.load()
    documents = add_metadata(documents, file_path)

    print(f"[Loader] ✅ {len(documents)} document(s) chargé(s) depuis '{Path(file_path).name}'")

    return documents


def load_directory(directory_path: str) -> List[Document]:
    """
    Charge tous les fichiers supportés d'un dossier.

    Args:
        directory_path : chemin vers le dossier

    Returns:
        List[Document] — tous les documents du dossier combinés

    Raises:
        FileNotFoundError : si le dossier n'existe pas

    Exemple:
        >>> docs = load_directory("data/csv_files/")
        >>> print(f"{len(docs)} documents chargés")
    """
    directory = Path(directory_path)

    if not directory.exists():
        raise FileNotFoundError(f"Dossier introuvable : {directory_path}")

    all_documents = []

    for file_path in sorted(directory.iterdir()):
        if not file_path.is_file():
            continue
        if file_path.suffix.lower() not in SUPPORTED_EXTENSIONS:
            continue

        try:
            docs = load_file(str(file_path))
            all_documents.extend(docs)
        except Exception as e:
            print(f"[Loader] ⚠️ Erreur sur '{file_path.name}' : {e}")
            continue

    print(f"[Loader] 📁 Total : {len(all_documents)} document(s) depuis '{directory.name}'")

    return all_documents